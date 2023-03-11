if __name__ == '__main__':
    __package__ = "presentations"
    pass
from .PresentationLogger import logger, PPTX_INFO, PPTX_DEBUG, PPTX_WARNING

import warnings
import itertools
import os
import manim
import pptx
import pptx.slide
import subprocess
import lxml.etree as etree
from functools import reduce

# noinspection PyProtectedMember
etreeElementClass = etree._Element

# Plan to convert to pptx in a way that can be used with teams:
#   Each "Slide" will create multiple pptx slides.
#   As a test, each pptx slide will hold a single animation, and those animations belonging to the same "Slide" will
#   automatically move on to the next slide, except for the last animation.
#   For example, a Slide with:
#           self.play(...)
#           self.wait(5)
#           self.play(...)
#   would result in three different slides, the first two of which automatically advance to the next slide when they
#   are complete.
#
#   Eventually, I should incorporate calls to "self.wait" to use the power point wait functionality, but only once
#   I've guaranteed that this strategy will allow me to use it with teams AND solve the issue I'm having with the
#   desktop version.
#       I also will need to make sure the images pause with correct displays (since currently this is being accounted
#       for with a call to self.wait).  It may be better to generate an image for each "self.wait" command and put that
#       up on the slide instead of a video.
#   If this works, each "Slide" should be a pptx section so they are grouped within power point.
#   I should also look into slide numbering, possible with <section#>.<slide#>
#   I should also look into the "Title" attributes of the slides, possibly having some sort of connection to the type
#   of animation.

# DONE: the above strategy will NOT be sufficient to get things working with microsoft teams presentations BECAUSE:
#  teams presentations do not support automatically advancing slides (a.k.a. automatic transitions)
#  ~~~
#  Instead, if I want to get it working with teams, I need to combine all the animations for a single slide into one
#  video, then put it in a single slide.  Use the manim combination tool, look in manim.scene.scene_file_writer.py:553

# TODO (2023-03-04 @ 09:10:37): I can get the slides to show their end state by putting an image of the end state in
#  the front of the slide (i.e. the next slide's start image) and then immediately removing the image once the slide
#  starts.  This might be jittery though, so only do this if necessary.
#  -> A better way to do this would be to "appear" the image in the front foreground AFTER the video completes (that
#     way the image is there when viewing the slide in the deck, but not there when the animation is playing)

# Possible future issues:
#   -There may still be issues with having to restart the slide show when presenting from desktop if the number of
#   animations passes ~60-70, I don't know if this can be resolved.  (might be a good idea to automatically include a
#   "restart" note in the notes when this threshold is passed)
#   -There may be issues if the animation for a single slide gets too long, and powerpoint then automatically reduces
#   the video quality.  (might be a good idea to identify this threshold and notify the user and/or automatically split
#   the slide into two animations when this happens)
#   -There may also still be issues with MS Teams because I haven't tested it rigorously (though it seems to work fine
#   from initial tests)


url_schema = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


# Quick and dirty implementation, copying most of PPTXScene from manim_pptx
class PPTXScene(manim.Scene):
    def __init__(self, *args, **kwargs):
        self.output_folder = kwargs.pop("output_folder", "./pptx/")
        self.temporary_dir = kwargs.pop("temporary_dir", "./temp/")
        super(PPTXScene, self).__init__(*args, **kwargs)

        self.slides = list()

        self.currentSlide = 1
        self.currentAnimation = 0
        self.currentSlideAnimations = 0
        self.slideStartAnimation = 0
        pass

    def construct(self, *args, **kwargs):
        return super(PPTXScene, self).construct()

    def play(self, *args, **kwargs):
        if self.currentSlideAnimations + 10 >= manim.config.max_files_cached:
            new_max_files = (3 * manim.config.max_files_cached // 2)
            logger.log(PPTX_WARNING,
                       f'The number of animations in slide {self.currentSlide} is approaching the max number of '
                       f'cached files.  Adjusting the max number of cached files from {manim.config.max_files_cached} '
                       f'to {new_max_files}.  To avoid this, manually adjust manim.config.max_files_cached in your '
                       f'script.')
            manim.config.max_files_cached = new_max_files
            pass
        super(PPTXScene, self).play(*args, **kwargs)
        self.currentAnimation += 1
        self.currentSlideAnimations += 1
        logger.log(PPTX_INFO, f"Add animation: {self.currentAnimation}")
        pass

    def wait(self, *args, **kwargs):
        super(PPTXScene, self).wait(*args, **kwargs)
        # self.currentAnimation += 1
        pass

    def all_endSlide(self, loop=False, autonext=False, notes=None, shownextnotes=False):
        warnings.warn("The 'all_endSlide' method is deprecated, "
                      "use 'endSlide' instead", DeprecationWarning, 2)
        logger.log(PPTX_INFO, f"End slide: {self.currentSlide} with animations "
                              f"[{self.slideStartAnimation},  {self.currentAnimation}]")
        self.slides.append(dict(
            type="loop" if loop else "slide",
            start=self.slideStartAnimation,
            end=self.currentAnimation,
            number=self.currentSlide,
            autonext=autonext,
            notes=notes,
            shownextnotes=shownextnotes,
        ))
        self.currentSlide += 1
        self.slideStartAnimation = self.currentAnimation
        self.currentSlideAnimations = 0
        pass

    def endSlide(self, loop=False, autonext=False, notes=None, shownextnotes=False):
        # TODO (2023-03-11 @ 09:47:34): Add a "quick=False" keyword argument that, when true, will, for the slide, set
        #  the _wait_override_ to a small value (~0.2), then set the wait time back to the original _wait_override_,
        #  and prepend a "(QUICK)" tag to the notes
        #  OR it might be better for this to go in the Slide.setup method
        logger.log(PPTX_INFO,
                   f"End slide: {self.currentSlide} with animations "
                   f"[{self.slideStartAnimation},  {self.currentAnimation}]")
        slide_info = dict(
            type="loop" if loop else "slide",
            start=self.slideStartAnimation,
            end=self.currentAnimation,
            number=self.currentSlide,
            autonext=autonext,
            notes=notes,
            shownextnotes=shownextnotes,
        )
        slide_info.update(
            self.__combineSlideAnimations__(self.currentSlide, self.slideStartAnimation, self.currentAnimation))
        self.slides.append(slide_info)

        self.currentSlide += 1
        self.slideStartAnimation = self.currentAnimation
        self.currentSlideAnimations = 0
        pass

    def __combineSlideAnimations__(self, slide_number: int, start_index: int, end_index: int):
        """Combine all animations for the slide into a single video."""
        slide_movie_files = self.renderer.file_writer.partial_movie_files[start_index:end_index]
        first_movie_file = slide_movie_files[0]
        # TODO (2023-03-05 @ 11:40:39): hash combined video (just hash combination of the contributing video hashes)
        #  so it can be reused (although this barely takes any time once the component videos are made, so maybe
        #  don't bother)
        slide_movie_file = os.path.join(self.temporary_dir,
                                        f"slide_{slide_number}" +
                                        os.path.splitext(os.path.basename(first_movie_file))[1])

        # This is very messy, and should eventually use more of the components seen in
        # manim.scene.scene_file_writer.py
        logger.log(PPTX_DEBUG, f"Combining animations [{start_index}, {end_index}]")
        self.renderer.file_writer.combine_files(slide_movie_files, slide_movie_file,
                                                create_gif=False, includes_sound=False)

        # Create thumbnail for animation
        thumbnail_file = os.path.join(
            self.temporary_dir,
            # os.path.splitext(os.path.basename(first_movie_file))[0] + ".png"
            f"slide_{slide_number}_thumbnail.png"
        )
        self.save_video_thumb(first_movie_file, thumbnail_file)

        return_data = {
            'slide_movie_file': slide_movie_file,
            'thumbnail_file': thumbnail_file,
        }
        return return_data

    @staticmethod
    def save_video_thumb(filename, image_name):
        subprocess.run([
            # constants.FFMPEG_BIN,  # This worked as of manim v 0.15.2
            # 'ffmpeg',  # hardcode option  (this is a workaround b/c manim.constants no longer has FFMPEG_BIN)
            manim.config.ffmpeg_executable,
            # This works as of manim v 0.17.2
            '-i', filename,
            '-vframes', '1',  # one frame
            '-loglevel', 'error',
            '-y',  # overwrite
            image_name,
        ], stdout=subprocess.PIPE)
        return image_name

    @staticmethod
    def get_dur(filename):
        return int(float(subprocess.check_output([
            "ffprobe",
            '-i', filename,
            "-show_entries", "format=duration",  # show duration
            "-v", "quiet",  # hide other outputs
            "-of", "csv=p=0",  # only number
        ]).decode("utf-8").strip()) * 1000)

    def old_render(self, *args, **kwargs):
        warnings.warn("The 'old_render' method is deprecated, "
                      "use 'render' instead", DeprecationWarning, 2)
        super(PPTXScene, self).render(*args, **kwargs)

        if not os.path.exists(self.output_folder):
            os.mkdir(self.output_folder)

        if not os.path.exists(self.temporary_dir):
            os.mkdir(self.temporary_dir)

        logger.log(PPTX_INFO, "Creating PPTX")

        prs = pptx.Presentation(pptx=os.path.join(os.path.split(__file__)[0], "template.pptx"))

        prs.slide_width = self.camera.pixel_width * 9525  # pixels to emu
        prs.slide_height = self.camera.pixel_height * 9525

        blank_slide_layout = prs.slide_layouts[6]

        for tslidei, tslide in enumerate(self.slides):
            logger.log(PPTX_DEBUG, f"Add slide {tslidei} with animations [{tslide['start']}, {tslide['end']}]")

            slide_movie_files = self.renderer.file_writer.partial_movie_files[tslide["start"]:tslide["end"]]

            slide = prs.slides.add_slide(blank_slide_layout)

            notes = tslide["notes"] if tslide["notes"] else ""

            if tslide["shownextnotes"] and len(self.slides) > tslidei + 1:
                notes += "\n" + "\n".join(list(map(lambda x: "> " + x, self.slides[tslidei + 1]["notes"].split("\n"))))

            slide.notes_slide.notes_text_frame.text = notes

            pics = list()

            for src_file in slide_movie_files:
                thumb_file = os.path.join(self.temporary_dir, os.path.basename(src_file) + ".png")
                self.save_video_thumb(src_file, thumb_file)

                logger.log(PPTX_DEBUG, f"adding video {src_file}")
                clip = slide.shapes.add_movie(src_file, 0, 0, prs.slide_width, prs.slide_height,
                                              mime_type='video/mp4', poster_frame_image=thumb_file)

                pics.append(dict(
                    id=clip.element[0][0].attrib.get("id"),
                    dur=self.get_dur(src_file),
                ))
                pass

            if len(pics) > 0:

                outerchildTnLst = slide.element[2][0][0][0][0]

                if tslide["autonext"]:
                    self.addAutoNext(slide)
                    outerchildTnLst = slide.element[3][0][0][0][0]

                seq = etree.Element(url_schema + "seq", concurrent="1", nextAc="seek")
                outerchildTnLst.insert(0, seq)

                cTnIDCounter = itertools.count(2)
                childTnLst = self.addCTn(cTnIDCounter, tslide, pics, seq)

                self.addPrevCondLst(seq)
                self.addNextCondLst(seq)

                currentdelay = 0
                for i, pic in enumerate(pics):
                    self.addToFrontEffect(cTnIDCounter, currentdelay, pic, i, childTnLst)
                    self.playEffect(cTnIDCounter, currentdelay, pic, childTnLst)
                    currentdelay += pic["dur"]
                    if i + 1 != len(pics):  # or tslide["type"] == "loop":
                        self.addToBackEffect(cTnIDCounter, currentdelay, pic, childTnLst)
                        pass
                    pass

                for i in range(1, len(outerchildTnLst)):
                    outerchildTnLst[i][0][0].attrib["id"] = str(next(cTnIDCounter))
                    pass
                pass
            pass

        presentation_name = getattr(type(self), "__presentation_name__", type(self).__name__)
        prs.save(os.path.join(self.output_folder, presentation_name + '.pptx'))
        pass

    def many_render(self, *args, **kwargs):
        warnings.warn("The 'many_render' method is deprecated, "
                      "use 'render' instead", DeprecationWarning, 2)
        # TODO (2023-03-03 @ 07:49:59): DOCUMENT!!!
        super(PPTXScene, self).render(*args, **kwargs)
        # TODO (2023-03-03 @ 11:04:49): remove "advance slide on mouse click" from pure animation slides
        if not os.path.exists(self.output_folder):
            os.mkdir(self.output_folder)
        if not os.path.exists(self.temporary_dir):
            os.mkdir(self.temporary_dir)

        logger.log(PPTX_INFO, "Creating PPTX")

        # open/load-in template presentation
        prs = pptx.Presentation(pptx=os.path.join(os.path.split(__file__)[0], "template.pptx"))

        prs.slide_width = self.camera.pixel_width * 9525  # pixels to emu
        prs.slide_height = self.camera.pixel_height * 9525

        blank_slide_layout = prs.slide_layouts[6]

        for tslidei, tslide in enumerate(self.slides):
            logger.log(PPTX_DEBUG, f"Add slide {tslidei} with animations [{tslide['start']}, {tslide['end']}]")

            slide_movie_files = self.renderer.file_writer.partial_movie_files[tslide["start"]:tslide["end"]]

            notes = tslide["notes"] if tslide["notes"] else ""

            if tslide["shownextnotes"] and len(self.slides) > tslidei + 1:
                notes += "\n" + "\n".join(list(map(lambda x: "> " + x, self.slides[tslidei + 1]["notes"].split("\n"))))

            # Loop over animations associated with the slide, making one slide for each animation and setting them to
            # automatically advance to the next slide, except for the last animation.
            n_animations = len(slide_movie_files)
            for offset_animation_index, src_file in enumerate(slide_movie_files, start=1 - n_animations):
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.notes_slide.notes_text_frame.text = notes + ("\n\n(ANIMATING...)"
                                                                   if offset_animation_index else "")

                thumb_file = os.path.join(self.temporary_dir,
                                          os.path.splitext(os.path.basename(src_file))[0] + ".png")
                self.save_video_thumb(src_file, thumb_file)
                logger.log(PPTX_DEBUG, f"adding video {src_file}")
                clip = slide.shapes.add_movie(src_file, 0, 0, prs.slide_width, prs.slide_height,
                                              mime_type='video/mp4', poster_frame_image=thumb_file)
                image_dict = {
                    "id": clip.element[0][0].attrib.get("id"),
                    "dur": PPTXScene.get_dur(src_file),
                }

                if offset_animation_index or tslide["autonext"]:
                    self.addAutoNext(slide)
                    outerchildTnLst = slide.element[3][0][0][0][0]
                else:
                    outerchildTnLst = slide.element[2][0][0][0][0]

                # Need to figure out what cTn means (or possibly childTn

                # No idea what this section is doing...
                #   this has something to do with storing the animation sequence
                seq = etree.Element(url_schema + "seq", concurrent="1", nextAc="seek")
                outerchildTnLst.insert(0, seq)

                cTnIDCounter = itertools.count(2)  # Some sort of counter
                if tslide["type"] == "loop":
                    tslide["type"] = "slide"
                    logger.warning("looping slides are not allowed for multi-rendered presentations. "
                                   "Slide type has been changed to non-looping slide.")
                    pass
                childTnLst = self.addCTn(cTnIDCounter, tslide, [], seq)
                self.addPrevCondLst(seq)  # Allows individual animation slide to be reviewed without automatically advancing
                self.addNextCondLst(seq)  # remove necessity for trigger for videos (animations)

                # Add effect to play the animation, starting immediately
                currentdelay = 0
                self.playEffect(cTnIDCounter, currentdelay, image_dict, childTnLst)

                # No idea what this part is doing...
                #   This makes it so the slide show goes smoothly onto the end-show slide (and possibly others)
                for i in range(1, len(outerchildTnLst)):
                    outerchildTnLst[i][0][0].attrib["id"] = str(next(cTnIDCounter))
                    pass
                pass

            pass

        presentation_name = getattr(type(self), "__presentation_name__", type(self).__name__)
        prs.save(os.path.join(self.output_folder, presentation_name + '.pptx'))
        pass

    def all_render(self, *args, **kwargs):
        warnings.warn("The 'all_render' method is deprecated, "
                      "use 'render' instead", DeprecationWarning, 2)
        presentation_name = getattr(type(self), "__presentation_name__", type(self).__name__) + '.pptx'
        presentation_path = os.path.join(self.output_folder, presentation_name)
        # Check if file already exists
        if os.path.isfile(presentation_path):
            # Check to make sure saving the resulting power point will be possible before spending time to create it
            try:
                os.remove(presentation_path)
                logger.log(PPTX_INFO, f'Removed previously existing PowerPoint file.')
                pass
            except PermissionError as e:
                e.strerror = f'Currently existing PowerPoint file is protected from modification by manim; the ' \
                             f'resulting presentation is not be able to be saved with the name {presentation_name}' \
                             f'.  Please manually rename or delete the existing file, then try again.' \
                             f'\n\tFile location:'
                raise
            pass

        # TODO (2023-03-03 @ 07:49:59): DOCUMENT!!!
        super(PPTXScene, self).render(*args, **kwargs)
        if not os.path.exists(self.output_folder):
            os.mkdir(self.output_folder)
        if not os.path.exists(self.temporary_dir):
            os.mkdir(self.temporary_dir)
        logger.log(PPTX_INFO, "Creating PPTX")

        # open/load-in template presentation
        prs = pptx.Presentation(pptx=os.path.join(os.path.split(__file__)[0], "template.pptx"))

        prs.slide_width = self.camera.pixel_width * 9525  # pixels to emu
        prs.slide_height = self.camera.pixel_height * 9525

        blank_slide_layout = prs.slide_layouts[6]

        for tslidei, tslide in enumerate(self.slides):
            logger.log(PPTX_DEBUG, f"Add slide {tslidei} with animations [{tslide['start']}, {tslide['end']}]")

            slide_movie_files = self.renderer.file_writer.partial_movie_files[tslide["start"]:tslide["end"]]
            notes = tslide["notes"] if tslide["notes"] else ""
            if tslide["shownextnotes"] and len(self.slides) > tslidei + 1:
                notes += "\n" + "\n".join(list(map(lambda x: "> " + x, self.slides[tslidei + 1]["notes"].split("\n"))))

            # This is very dirty, and should eventually use more of the components seen in
            # manim.scene.scene_file_writer.py
            first_movie_file = slide_movie_files[0]
            slide_movie_file = os.path.join(self.temporary_dir,
                                            f"slide_{tslidei}" +
                                            os.path.splitext(os.path.basename(first_movie_file))[1])
            logger.log(PPTX_DEBUG, f"Combine animations [{tslide['start']}, {tslide['end']}]")
            self.renderer.file_writer.combine_files(slide_movie_files, slide_movie_file,
                                                    create_gif=False, includes_sound=False)

            slide = prs.slides.add_slide(blank_slide_layout)
            slide.notes_slide.notes_text_frame.text = notes
            thumbnail_file = os.path.join(self.temporary_dir,
                                          # os.path.splitext(os.path.basename(first_movie_file))[0] + ".png"
                                          f"slide_{tslidei}_thumbnail.png"
                                          )
            self.save_video_thumb(first_movie_file, thumbnail_file)
            logger.log(PPTX_DEBUG, f"adding animation {slide_movie_file}")
            clip = slide.shapes.add_movie(slide_movie_file, 0, 0, prs.slide_width, prs.slide_height,
                                          mime_type='video/mp4', poster_frame_image=thumbnail_file)
            image_dict = {
                "id": clip.element[0][0].attrib.get("id"),
                "dur": PPTXScene.get_dur(slide_movie_file),
            }
            if tslide["autonext"]:
                self.addAutoNext(slide)
                outerchildTnLst = slide.element[3][0][0][0][0]
                pass
            else:
                outerchildTnLst = slide.element[2][0][0][0][0]
                pass

            # Need to figure out what cTn means (or possibly childTn)

            # No idea what this section is doing...
            #   I think this has something to do with storing the animation sequence
            seq = etree.Element(url_schema + "seq", concurrent="1", nextAc="seek")
            outerchildTnLst.insert(0, seq)
            cTnIDCounter = itertools.count(2)  # Some sort of counter
            childTnLst = self.addCTn(cTnIDCounter, tslide, [image_dict], seq)
            self.addPrevCondLst(seq)  # Allows individual animation slide to be reviewed without automatically advancing
            self.addNextCondLst(seq)  # remove necessity for trigger for videos (animations)

            # Add effect to play the animation, starting immediately
            currentdelay = 0
            self.playEffect(cTnIDCounter, currentdelay, image_dict, childTnLst)

            # No idea what this part is doing...
            #   I think this makes it so the slide show goes smoothly onto the end-show slide (and possibly others)
            for i in range(1, len(outerchildTnLst)):
                outerchildTnLst[i][0][0].attrib["id"] = str(next(cTnIDCounter))
                pass
            pass

        # Save presentation
        prs.save(presentation_path)
        logger.log(PPTX_INFO, f'PowerPoint written to: {presentation_name} in {self.output_folder}')
        # print(self.renderer.file_writer.partial_movie_directory)
        pass

    def render(self, *args, **kwargs):
        presentation_name = getattr(type(self), "__presentation_name__", type(self).__name__) + '.pptx'
        presentation_path = os.path.join(self.output_folder, presentation_name)
        # Check if file already exists
        if os.path.isfile(presentation_path):
            # Check to make sure saving the resulting power point will be possible before spending time to create it
            try:
                os.remove(presentation_path)
                logger.log(PPTX_INFO, f'Removed previously existing PowerPoint file.')
                pass
            except PermissionError as e:
                e.strerror = f'Currently existing PowerPoint file is protected from modification by manim; the ' \
                             f'resulting presentation is not be able to be saved with the name {presentation_name}' \
                             f'.  Please manually rename or delete the existing file, then try again.' \
                             f'\n\tFile location:'
                raise
            pass

        # TODO (2023-03-03 @ 07:49:59): DOCUMENT!!!
        if not os.path.exists(self.output_folder):
            os.mkdir(self.output_folder)
        if not os.path.exists(self.temporary_dir):
            os.mkdir(self.temporary_dir)
        super(PPTXScene, self).render(*args, **kwargs)
        logger.log(PPTX_INFO, "Creating PPTX")

        # open/load-in template presentation
        prs = pptx.Presentation(pptx=os.path.join(os.path.split(__file__)[0], "template.pptx"))
        prs.slide_width = self.camera.pixel_width * 9525  # pixels to emu
        prs.slide_height = self.camera.pixel_height * 9525
        blank_slide_layout = prs.slide_layouts[6]

        for slide_index, tslide in enumerate(self.slides):
            logger.log(PPTX_DEBUG,
                       f"Add slide {tslide['number']} with animations [{tslide['start']}, {tslide['end']}]")

            notes = tslide["notes"] if tslide["notes"] else ""
            if tslide["shownextnotes"] and len(self.slides) > slide_index + 1:
                notes += "\n" + "\n".join(list(map(
                    lambda x: "> " + x, self.slides[slide_index + 1]["notes"].split("\n"))))
                pass

            slide = prs.slides.add_slide(blank_slide_layout)
            slide.notes_slide.notes_text_frame.text = notes
            slide_movie_file = tslide['slide_movie_file']
            thumbnail_file = tslide['thumbnail_file']
            logger.log(PPTX_DEBUG, f"adding animation {slide_movie_file}")
            clip = slide.shapes.add_movie(slide_movie_file, 0, 0, prs.slide_width, prs.slide_height,
                                          mime_type='video/mp4', poster_frame_image=thumbnail_file)
            image_dict = {
                "id": clip.element[0][0].attrib.get("id"),
                "dur": PPTXScene.get_dur(slide_movie_file),
            }
            if tslide["autonext"]:
                self.addAutoNext(slide)
                outerchildTnLst = slide.element[3][0][0][0][0]
                pass
            else:
                outerchildTnLst = slide.element[2][0][0][0][0]
                pass

            # Need to figure out what cTn means (or possibly childTn)

            # No idea what this section is doing...
            #   I think this has something to do with storing the animation sequence
            seq = etree.Element(url_schema + "seq", concurrent="1", nextAc="seek")
            outerchildTnLst.insert(0, seq)
            cTnIDCounter = itertools.count(2)  # Some sort of counter
            childTnLst = self.addCTn(cTnIDCounter, tslide, [image_dict], seq)
            self.addPrevCondLst(seq)  # Allows individual animation slide to be reviewed without automatically advancing
            self.addNextCondLst(seq)  # remove necessity for trigger for videos (animations)

            # Add effect to play the animation, starting immediately
            currentdelay = 0
            self.playEffect(cTnIDCounter, currentdelay, image_dict, childTnLst)

            # No idea what this part is doing...
            #   I think this makes it so the slide show goes smoothly onto the end-show slide (and possibly others)
            for i in range(1, len(outerchildTnLst)):
                outerchildTnLst[i][0][0].attrib["id"] = str(next(cTnIDCounter))
                pass
            pass

        # Save presentation
        prs.save(presentation_path)
        logger.log(PPTX_INFO, f'PowerPoint written to: {presentation_name} in {self.output_folder}')
        # print(self.renderer.file_writer.partial_movie_directory)
        pass

    # static methods for controlling slide animations, transitions, and other components.
    #   see https://python-pptx.readthedocs.io/en/latest/dev/analysis/shp-movie.html for more information

    # TODO (2023-03-06 @ 10:30:19): Convert these to a class of their own (this should make it more OOP)

    @staticmethod
    def addAutoNext(slide: pptx.slide.Slide):
        transition = etree.Element(url_schema + "transition", {
            "spd": "slow",
            "advTm": "0",
        })
        slide.element.insert(2, transition)
        pass

    @staticmethod
    def addCTn(cTnIDCounter: itertools.count, tslide: dict, pics: list, seq: etreeElementClass):
        innercTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), dur="indefinite", nodeType="mainSeq")
        childTnLst = etree.Element(url_schema + "childTnLst")
        par1 = etree.Element(url_schema + "par")
        cTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), fill="hold")
        if tslide["type"] == "loop":
            cTn.attrib["dur"] = str(reduce(lambda x, y: x + y, [p["dur"] for p in pics]))
            cTn.attrib["repeatCount"] = "indefinite"
        stCondLst = etree.Element(url_schema + "stCondLst")
        cond1 = etree.Element(url_schema + "cond", delay="indefinite")
        cond2 = etree.Element(url_schema + "cond", evt="onBegin", delay="0")
        cond2tn = etree.Element(url_schema + "tn", val="2")
        cond2.append(cond2tn)
        stCondLst.append(cond1)
        stCondLst.append(cond2)
        cTn.append(stCondLst)
        childTnLst2 = etree.Element(url_schema + "childTnLst")
        cTn.append(childTnLst2)

        par1.append(cTn)
        childTnLst.append(par1)
        innercTn.append(childTnLst)
        seq.append(innercTn)
        return childTnLst2

    @staticmethod
    def addPrevCondLst(seq: etreeElementClass):
        prevCondLst = etree.Element(url_schema + "prevCondLst")
        cond = etree.Element(url_schema + "cond", evt="onPrev", delay="0")
        tgtEl = etree.Element(url_schema + "tgtEl")
        sldTgt = etree.Element(url_schema + "sldTgt")
        tgtEl.append(sldTgt)
        cond.append(tgtEl)
        prevCondLst.append(cond)
        seq.append(prevCondLst)
        pass

    @staticmethod
    def addNextCondLst(seq: etreeElementClass):
        nextCondLst = etree.Element(url_schema + "nextCondLst")
        cond = etree.Element(url_schema + "cond", evt="onNext", delay="0")
        tgtEl = etree.Element(url_schema + "tgtEl")
        sldTgt = etree.Element(url_schema + "sldTgt")
        tgtEl.append(sldTgt)
        cond.append(tgtEl)
        nextCondLst.append(cond)
        seq.append(nextCondLst)
        pass

    @staticmethod
    def addToFrontEffect(cTnIDCounter: itertools.count, currentdelay: int, image_dict: dict, i: int,
                         childTnLst: etreeElementClass):
        par = etree.Element(url_schema + "par")
        cTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), fill="hold")
        stCondLst = etree.Element(url_schema + "stCondLst")
        cond = etree.Element(url_schema + "cond", delay=str(currentdelay))
        stCondLst.append(cond)
        cTn.append(stCondLst)

        innerchildTnLst = etree.Element(url_schema + "childTnLst")
        innerPar = etree.Element(url_schema + "par")
        innercTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), presetID="1", presetClass="entr",
                                 presetSubtype="0", fill="hold", nodeType="withEffect" if i == 0 else "afterEffect")
        innerstCondLst = etree.Element(url_schema + "stCondLst")
        innercond = etree.Element(url_schema + "cond", delay="0")
        innerstCondLst.append(innercond)
        innercTn.append(innerstCondLst)

        innerInnerChildTnLst = etree.Element(url_schema + "childTnLst")
        innercTn.append(innerInnerChildTnLst)

        setElement = etree.Element(url_schema + "set")
        cBhvr = etree.Element(url_schema + "cBhvr")

        cBhvrcTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), dur="1", fill="hold")
        cBhvrcTnstCondLst = etree.Element(url_schema + "stCondLst")
        cBhvrcTncond = etree.Element(url_schema + "cond", delay="0")
        cBhvrcTnstCondLst.append(cBhvrcTncond)
        cBhvrcTn.append(cBhvrcTnstCondLst)
        cBhvr.append(cBhvrcTn)
        tgtEl = etree.Element(url_schema + "tgtEl")
        spTgt = etree.Element(url_schema + "spTgt", spid=str(image_dict["id"]))
        tgtEl.append(spTgt)
        cBhvr.append(tgtEl)
        attrNameLst = etree.Element(url_schema + "attrNameLst")
        attrName = etree.Element(url_schema + "attrName")
        attrName.text = "style.visibility"
        attrNameLst.append(attrName)
        cBhvr.append(attrNameLst)

        setElement.append(cBhvr)
        to = etree.Element(url_schema + "to")
        strVal = etree.Element(url_schema + "strVal", val="visible")
        to.append(strVal)
        setElement.append(to)

        innerInnerChildTnLst.append(setElement)

        innerPar.append(innercTn)
        innerchildTnLst.append(innerPar)
        cTn.append(innerchildTnLst)

        par.append(cTn)
        childTnLst.append(par)
        pass

    @staticmethod
    def addToBackEffect(cTnIDCounter: itertools.count, currentdelay: int, image_dict: dict,
                        childTnLst: etreeElementClass):
        par = etree.Element(url_schema + "par")
        cTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), fill="hold")
        stCondLst = etree.Element(url_schema + "stCondLst")
        cond = etree.Element(url_schema + "cond", delay=str(currentdelay))
        stCondLst.append(cond)
        cTn.append(stCondLst)

        innerchildTnLst = etree.Element(url_schema + "childTnLst")
        innerPar = etree.Element(url_schema + "par")
        innercTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), presetID="1", presetClass="exit",
                                 presetSubtype="0", fill="hold", nodeType="afterEffect")
        innerstCondLst = etree.Element(url_schema + "stCondLst")
        innercond = etree.Element(url_schema + "cond", delay="0")
        innerstCondLst.append(innercond)
        innercTn.append(innerstCondLst)

        innerInnerChildTnLst = etree.Element(url_schema + "childTnLst")
        innercTn.append(innerInnerChildTnLst)

        setElement = etree.Element(url_schema + "set")
        cBhvr = etree.Element(url_schema + "cBhvr")

        cBhvrcTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), dur="1", fill="hold")
        cBhvrcTnstCondLst = etree.Element(url_schema + "stCondLst")
        cBhvrcTncond = etree.Element(url_schema + "cond", delay="0")
        cBhvrcTnstCondLst.append(cBhvrcTncond)
        cBhvrcTn.append(cBhvrcTnstCondLst)
        cBhvr.append(cBhvrcTn)
        tgtEl = etree.Element(url_schema + "tgtEl")
        spTgt = etree.Element(url_schema + "spTgt", spid=str(image_dict["id"]))
        tgtEl.append(spTgt)
        cBhvr.append(tgtEl)
        attrNameLst = etree.Element(url_schema + "attrNameLst")
        attrName = etree.Element(url_schema + "attrName")
        attrName.text = "style.visibility"
        attrNameLst.append(attrName)
        cBhvr.append(attrNameLst)

        setElement.append(cBhvr)
        to = etree.Element(url_schema + "to")
        strVal = etree.Element(url_schema + "strVal", val="hidden")
        to.append(strVal)
        setElement.append(to)

        innerInnerChildTnLst.append(setElement)

        cmd = etree.Element(url_schema + "cmd", type="call", cmd="stop")
        cBhvr = etree.Element(url_schema + "cBhvr")
        cBhvrcTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), dur="1", fill="hold")
        cBhvrcTnstCondLst = etree.Element(url_schema + "stCondLst")
        cBhvrcTnstCondLstCond = etree.Element(url_schema + "cond", delay="0")
        cBhvrcTnstCondLst.append(cBhvrcTnstCondLstCond)
        cBhvrcTn.append(cBhvrcTnstCondLst)
        cBhvr.append(cBhvrcTn)
        tgtEl = etree.Element(url_schema + "tgtEl")
        spTgt = etree.Element(url_schema + "spTgt", spid=str(image_dict["id"]))
        tgtEl.append(spTgt)
        cBhvr.append(tgtEl)
        cmd.append(cBhvr)

        innerInnerChildTnLst.append(cmd)

        innerPar.append(innercTn)
        innerchildTnLst.append(innerPar)
        cTn.append(innerchildTnLst)

        par.append(cTn)
        childTnLst.append(par)
        pass

    @staticmethod
    def playEffect(cTnIDCounter: itertools.count, currentdelay: int, image_dict: dict,
                   childTnLst: etreeElementClass):
        par = etree.Element(url_schema + "par")
        cTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), fill="hold")
        stCondLst = etree.Element(url_schema + "stCondLst")
        cond = etree.Element(url_schema + "cond", delay=str(currentdelay))
        stCondLst.append(cond)
        cTn.append(stCondLst)

        innerchildTnLst = etree.Element(url_schema + "childTnLst")
        innerPar = etree.Element(url_schema + "par")
        innercTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), presetID="1", presetClass="mediacall",
                                 presetSubtype="0", fill="hold", nodeType="afterEffect")
        innerstCondLst = etree.Element(url_schema + "stCondLst")
        innercond = etree.Element(url_schema + "cond", delay="0")
        innerstCondLst.append(innercond)
        innercTn.append(innerstCondLst)

        innerInnerChildTnLst = etree.Element(url_schema + "childTnLst")
        innercTn.append(innerInnerChildTnLst)

        cmd = etree.Element(url_schema + "cmd", type="call", cmd="playFrom(0.0)")
        cBhvr = etree.Element(url_schema + "cBhvr")
        cBhvrcTn = etree.Element(url_schema + "cTn", id=str(next(cTnIDCounter)), dur=str(image_dict["dur"]),
                                 fill="hold")
        cBhvr.append(cBhvrcTn)
        tgtEl = etree.Element(url_schema + "tgtEl")
        spTgt = etree.Element(url_schema + "spTgt", spid=str(image_dict["id"]))
        tgtEl.append(spTgt)
        cBhvr.append(tgtEl)
        cmd.append(cBhvr)

        innerInnerChildTnLst.append(cmd)

        innerPar.append(innercTn)
        innerchildTnLst.append(innerPar)
        cTn.append(innerchildTnLst)

        par.append(cTn)
        childTnLst.append(par)
        pass

    pass


if __name__ == '__main__':
    """Testing for creating multiple power point slides for a single Slide decorated function."""

    # open/load in template presentation
    prs_t = pptx.Presentation(pptx=os.path.join(os.path.split(__file__)[0], "template.pptx"))
    blank_slide_layout_t = prs_t.slide_layouts[6]
    source_folder_t = "C:\\Users\\james\\AppData\\Local\\Programs\\Python\\Python38\\Lib\\site-packages" \
                      "\\presentations\\TestPartialVideos\\"
    slide_movie_files_t = [
        "274514146_2271961971_223132457.mp4",
        "207390714_4070753449_2157517633.mp4",
        "207390714_4075652138_2859536990.mp4",
        "207390714_3639972290_3722111848.mp4"
    ]
    slide_movie_files_t = list(map(source_folder_t.__add__, slide_movie_files_t))
    notes_t = "These are test notes."
    temporary_dir_t = "C:\\Users\\james\\AppData\\Local\\Programs\\Python\\Python38\\Lib\\site-packages" \
                      "\\presentations\\TestTempDir\\"
    output_folder_t = "C:\\Users\\james\\AppData\\Local\\Programs\\Python\\Python38\\Lib\\site-packages" \
                      "\\presentations\\TestOutput\\"
    if not os.path.exists(output_folder_t):
        os.mkdir(output_folder_t)
        pass
    if not os.path.exists(temporary_dir_t):
        os.mkdir(temporary_dir_t)
        pass

    n_animations_t = len(slide_movie_files_t)
    for offset_animation_index_t, src_file_t in enumerate(slide_movie_files_t, start=1 - n_animations_t):
        slide_t = prs_t.slides.add_slide(blank_slide_layout_t)
        slide_t.notes_slide.notes_text_frame.text = notes_t

        thumb_file_t = os.path.join(temporary_dir_t,
                                    os.path.splitext(os.path.basename(src_file_t))[0] + ".png")
        PPTXScene.save_video_thumb(src_file_t, thumb_file_t)
        clip_t = slide_t.shapes.add_movie(src_file_t, 0, 0, prs_t.slide_width, prs_t.slide_height,
                                          mime_type='video/mp4', poster_frame_image=thumb_file_t)
        image_dict_t = {
            "id": clip_t.element[0][0].attrib.get("id"),
            "dur": PPTXScene.get_dur(src_file_t),
        }

        if offset_animation_index_t:  # or tslide["autonext"]:
            PPTXScene.addAutoNext(slide_t)
            outerchildTnLst_t = slide_t.element[3][0][0][0][0]
        else:
            outerchildTnLst_t = slide_t.element[2][0][0][0][0]

        # Need to figure out what cTn means (or possibly childTn

        # No idea what this section is doing...
        #   this has something to do with storing the animation sequence
        seq_t = etree.Element(url_schema + "seq", concurrent="1", nextAc="seek")
        outerchildTnLst_t.insert(0, seq_t)

        cTnIDCounter_t = itertools.count(2)  # Some sort of counter
        tslide_t = {"type": "slide"}  # Assign the slide for this test
        childTnLst_t = PPTXScene.addCTn(cTnIDCounter_t, tslide_t, [image_dict_t], seq_t)
        PPTXScene.addPrevCondLst(seq_t)  # Allows individual animation slide review without automatically advancing
        PPTXScene.addNextCondLst(seq_t)  # Removes necessity for a trigger for videos (animations)

        # Add effect to play the animation, starting immediately
        currentdelay_t = 0
        PPTXScene.playEffect(cTnIDCounter_t, currentdelay_t, image_dict_t, childTnLst_t)

        # No idea what this part is doing...
        #   This makes it so the slide show goes smoothly onto the end-show slide (and possibly others)
        for i_t in range(1, len(outerchildTnLst_t)):
            outerchildTnLst_t[i_t][0][0].attrib["id"] = str(next(cTnIDCounter_t))
            pass
        pass

    prs_t.save(os.path.join(output_folder_t, 'test.pptx'))
    pass
